#!/usr/bin/env python3
"""导入任务支持文件类型冒烟测试。

脚本会生成每个声明支持后缀的最小样本，经 Java 本地小平台入口
POST /api/import-tasks 上传，并断言任务能进入可校验状态且至少生成 1 道题。
"""

from __future__ import annotations

import json
import mimetypes
import os
import shutil
import subprocess
import sys
import tempfile
import time
import urllib.error
import urllib.request
import uuid
from pathlib import Path
from typing import Any, Callable


PROJECT_ROOT = Path(__file__).resolve().parents[1]
WORKER_PYTHON = PROJECT_ROOT / "backend" / "python-worker" / ".venv" / "bin" / "python"
BASE_URL = os.environ.get("AI_GENERATION_BASE_URL", "http://localhost:8018").rstrip("/")
SUPPORTED_ORDER = [
    ".md",
    ".markdown",
    ".pdf",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".tif",
    ".tiff",
    ".doc",
    ".docx",
    ".pptx",
    ".xlsx",
]
QUESTION_TEXT = "1. x + 1 = 2, find x.\nA. 0\nB. 1\nC. 2\nD. 3\n"


def ensure_worker_python() -> None:
    if not WORKER_PYTHON.exists():
        return
    if Path(sys.executable) != WORKER_PYTHON:
        os.execv(str(WORKER_PYTHON), [str(WORKER_PYTHON), *sys.argv])


def request(method: str, path: str, payload: bytes | None = None, headers: dict[str, str] | None = None, timeout: int = 90) -> Any:
    req = urllib.request.Request(BASE_URL + path, data=payload, headers=headers or {}, method=method)
    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            raw = resp.read()
            content_type = resp.headers.get("content-type", "")
            if "application/json" in content_type:
                return json.loads(raw.decode("utf-8") or "{}")
            return raw
    except urllib.error.HTTPError as exc:
        raw = exc.read().decode("utf-8", errors="replace")
        raise RuntimeError(f"{method} {path} -> HTTP {exc.code}: {raw[:800]}") from exc


def multipart(fields: dict[str, str], files: dict[str, Path]) -> tuple[bytes, str]:
    boundary = "----question-engine-filetype-" + uuid.uuid4().hex
    chunks: list[bytes] = []
    for name, value in fields.items():
        chunks.append(f"--{boundary}\r\n".encode())
        chunks.append(f'Content-Disposition: form-data; name="{name}"\r\n\r\n'.encode())
        chunks.append(value.encode("utf-8"))
        chunks.append(b"\r\n")
    for name, path in files.items():
        content_type = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
        chunks.append(f"--{boundary}\r\n".encode())
        chunks.append(f'Content-Disposition: form-data; name="{name}"; filename="{path.name}"\r\n'.encode())
        chunks.append(f"Content-Type: {content_type}\r\n\r\n".encode())
        chunks.append(path.read_bytes())
        chunks.append(b"\r\n")
    chunks.append(f"--{boundary}--\r\n".encode())
    return b"".join(chunks), f"multipart/form-data; boundary={boundary}"


def make_markdown(path: Path) -> None:
    path.write_text("# 文件类型测试\n\n" + QUESTION_TEXT, encoding="utf-8")


def make_pdf(path: Path) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    pdf = canvas.Canvas(str(path), pagesize=A4)
    pdf.setFont("Helvetica", 18)
    y = 780
    for line in QUESTION_TEXT.splitlines():
        pdf.drawString(72, y, line)
        y -= 34
    pdf.save()


def make_image(path: Path, image_format: str) -> None:
    from PIL import Image, ImageDraw, ImageFont

    image = Image.new("RGB", (1400, 900), "white")
    draw = ImageDraw.Draw(image)
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Supplemental/Arial.ttf", 56)
    except Exception:
        font = ImageFont.load_default()
    y = 90
    for line in QUESTION_TEXT.splitlines():
        draw.text((80, y), line, fill="black", font=font)
        y += 90
    image.save(path, format=image_format)


def make_doc(path: Path) -> None:
    rtf = path.with_suffix(".rtf")
    rtf.write_text(
        r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Arial;}} \f0\fs28 "
        r"1. x + 1 = 2, find x.\line A. 0   B. 1   C. 2   D. 3}",
        encoding="utf-8",
    )
    attempts: list[str] = []
    textutil = shutil.which("textutil")
    if textutil:
        result = subprocess.run(
            [textutil, "-convert", "doc", "-output", str(path), str(rtf)],
            capture_output=True,
            text=True,
            timeout=30,
        )
        if result.returncode == 0 and path.is_file() and path.stat().st_size > 0:
            return
        attempts.append(f"textutil exited with code {result.returncode}: {(result.stderr or result.stdout)[-500:]}")

    office = shutil.which("soffice") or shutil.which("libreoffice")
    if office:
        result = subprocess.run(
            [office, "--headless", "--convert-to", "doc", "--outdir", str(path.parent), str(rtf)],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode == 0 and path.is_file() and path.stat().st_size > 0:
            return
        attempts.append(f"soffice exited with code {result.returncode}: {(result.stderr or result.stdout)[-500:]}")

    detail = "; ".join(attempts) or "textutil and LibreOffice/soffice are unavailable"
    raise RuntimeError(f"unable to generate a real .doc smoke sample: {detail}")


def make_docx(path: Path) -> None:
    from docx import Document

    document = Document()
    document.add_heading("文件类型测试", level=1)
    for line in QUESTION_TEXT.splitlines():
        document.add_paragraph(line)
    document.save(path)


def make_pptx(path: Path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    box = slide.shapes.add_textbox(500000, 500000, 8000000, 3000000)
    frame = box.text_frame
    frame.text = "1. x + 1 = 2, find x."
    paragraph = frame.add_paragraph()
    paragraph.text = "A. 0   B. 1   C. 2   D. 3"
    presentation.save(path)


def make_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet["A1"] = "1. x + 1 = 2, find x."
    sheet["A2"] = "A. 0   B. 1   C. 2   D. 3"
    workbook.save(path)


def generators() -> dict[str, Callable[[Path], None]]:
    return {
        ".md": make_markdown,
        ".markdown": make_markdown,
        ".pdf": make_pdf,
        ".png": lambda path: make_image(path, "PNG"),
        ".jpg": lambda path: make_image(path, "JPEG"),
        ".jpeg": lambda path: make_image(path, "JPEG"),
        ".webp": lambda path: make_image(path, "WEBP"),
        ".tif": lambda path: make_image(path, "TIFF"),
        ".tiff": lambda path: make_image(path, "TIFF"),
        ".doc": make_doc,
        ".docx": make_docx,
        ".pptx": make_pptx,
        ".xlsx": make_xlsx,
    }


def runtime_extensions() -> list[str]:
    runtime = request("GET", "/api/capabilities/ocr-flow/runtime", timeout=15)
    return [str(item).lower() for item in runtime.get("allowedExtensions") or []]


def wait_import_task(task_id: str) -> dict[str, Any]:
    deadline = time.time() + int(os.environ.get("AI_GENERATION_FILETYPE_TIMEOUT_SECONDS", "240"))
    last: dict[str, Any] = {}
    while time.time() < deadline:
        last = request("GET", f"/api/import-tasks/{task_id}", timeout=60)
        if str(last.get("status") or "") not in {"处理中", "processing", ""}:
            return last
        time.sleep(2)
    return last


def test_extension(extension: str, workdir: Path) -> dict[str, Any]:
    generator = generators().get(extension)
    if generator is None:
        return {"extension": extension, "ok": False, "error": "没有对应的样本生成器"}

    sample = workdir / f"sample{extension}"
    generator(sample)
    body, content_type = multipart(
        {
            "stage": "高中",
            "subject": "数学",
            "grade": "高一",
            "year": "2026",
            "title": "FileTypeSmoke_" + extension.strip(".") + "_" + uuid.uuid4().hex[:8],
        },
        {"paperFile": sample},
    )
    task = request("POST", "/api/import-tasks", body, {"Content-Type": content_type}, timeout=90)
    task_id = str(task["id"])
    detail = wait_import_task(task_id)
    question_count = len(detail.get("questions") or [])
    ok = detail.get("status") in {"待校验", "部分完成", "已完成"} and question_count > 0
    result = {
        "extension": extension,
        "ok": ok,
        "taskId": task_id,
        "status": detail.get("status"),
        "paperOcrStatus": detail.get("paperOcrStatus"),
        "questionCount": question_count,
        "failureReason": detail.get("failureReason") or "",
    }
    if ok:
        request("DELETE", f"/api/import-tasks/{task_id}", timeout=30)
    return result


def main() -> None:
    ensure_worker_python()
    extensions = runtime_extensions()
    ordered_extensions = [item for item in SUPPORTED_ORDER if item in extensions]
    missing_from_order = [item for item in extensions if item not in SUPPORTED_ORDER]
    if missing_from_order:
        raise AssertionError(f"运行时声明了未纳入测试的后缀: {missing_from_order}")

    workdir = Path(tempfile.mkdtemp(prefix="question-engine-filetypes-"))
    failures: list[dict[str, Any]] = []
    for extension in ordered_extensions:
        result = test_extension(extension, workdir)
        print(json.dumps(result, ensure_ascii=False), flush=True)
        if not result["ok"]:
            failures.append(result)
    if failures:
        raise AssertionError("存在未跑通的支持文件类型: " + json.dumps(failures, ensure_ascii=False))
    print("所有声明支持的导入文件类型均已跑通")


if __name__ == "__main__":
    main()
